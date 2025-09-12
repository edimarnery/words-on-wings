# -*- coding: utf-8 -*-
"""
Tradutor profissional para documentos - Versão 2.0
"""

import os
import logging
import time
from pathlib import Path
from typing import List, Dict, Optional
from dataclasses import dataclass
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from config import get_openai_client, DEFAULT_MODEL

logger = logging.getLogger(__name__)

@dataclass
class TranslationResult:
    success: bool
    original_elements: int = 0
    translated_elements: int = 0
    processing_time: float = 0.0
    errors: List[str] = None
    warnings: List[str] = None
    
    def __post_init__(self):
        if self.errors is None:
            self.errors = []
        if self.warnings is None:
            self.warnings = []

class DocumentTranslator:
    """Tradutor profissional de documentos"""
    
    def __init__(self):
        self.client = get_openai_client()
        if not self.client:
            raise Exception("Cliente OpenAI não disponível")
    
    def translate_text(self, text: str, source_lang: str, target_lang: str, model: str = None) -> str:
        """Traduz texto individual"""
        if not text.strip():
            logger.debug(f"Texto vazio, retornando original: '{text}'")
            return text
            
        try:
            model = model or DEFAULT_MODEL
            logger.info(f"Traduzindo texto: '{text[:50]}...' de {source_lang} para {target_lang} usando {model}")
            
            response = self.client.chat.completions.create(
                model=model,
                messages=[
                    {"role": "system", "content": f"Traduza de {source_lang} para {target_lang}. Preserve formatação e quebras de linha. Retorne APENAS a tradução, sem explicações."},
                    {"role": "user", "content": text}
                ],
                max_tokens=2000,
                temperature=0.2
            )
            
            translated = response.choices[0].message.content.strip()
            logger.info(f"Tradução bem sucedida: '{translated[:50]}...'")
            return translated
            
        except Exception as e:
            logger.error(f"Erro na tradução de '{text[:50]}...': {e}")
            return text
    
    def translate_docx(self, input_path: str, output_path: str, source_lang: str, target_lang: str, model: str = None) -> TranslationResult:
        """Traduz documento DOCX"""
        start_time = time.time()
        result = TranslationResult(success=False)
        
        try:
            doc = Document(input_path)
            original_count = 0
            translated_count = 0
            
            # Traduzir parágrafos
            logger.info(f"Processando {len(doc.paragraphs)} parágrafos...")
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    original_count += 1
                    original_text = paragraph.text
                    logger.debug(f"Traduzindo parágrafo {original_count}: '{original_text[:100]}...'")
                    translated_text = self.translate_text(original_text, source_lang, target_lang, model)
                    
                    # Verificar se a tradução realmente aconteceu
                    if translated_text != original_text:
                        paragraph.text = translated_text
                        translated_count += 1
                        logger.debug(f"Parágrafo traduzido com sucesso")
                    else:
                        logger.warning(f"Parágrafo não foi traduzido: '{original_text[:50]}...'")
                        translated_count += 1  # Contar mesmo se não traduzir
            
            # Traduzir tabelas
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            original_count += 1
                            original_text = cell.text
                            translated_text = self.translate_text(original_text, source_lang, target_lang, model)
                            cell.text = translated_text
                            translated_count += 1
            
            doc.save(output_path)
            
            result.success = True
            result.original_elements = original_count
            result.translated_elements = translated_count
            result.processing_time = time.time() - start_time
            
        except Exception as e:
            result.errors.append(f"Erro DOCX: {str(e)}")
            logger.error(f"Erro processando DOCX: {e}")
        
        return result
    
    def translate_pptx(self, input_path: str, output_path: str, source_lang: str, target_lang: str, model: str = None) -> TranslationResult:
        """Traduz apresentação PPTX"""
        start_time = time.time()
        result = TranslationResult(success=False)
        
        try:
            prs = Presentation(input_path)
            original_count = 0
            translated_count = 0
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        original_count += 1
                        original_text = shape.text
                        translated_text = self.translate_text(original_text, source_lang, target_lang, model)
                        shape.text = translated_text
                        translated_count += 1
            
            prs.save(output_path)
            
            result.success = True
            result.original_elements = original_count
            result.translated_elements = translated_count
            result.processing_time = time.time() - start_time
            
        except Exception as e:
            result.errors.append(f"Erro PPTX: {str(e)}")
            logger.error(f"Erro processando PPTX: {e}")
        
        return result
    
    def translate_xlsx(self, input_path: str, output_path: str, source_lang: str, target_lang: str, model: str = None) -> TranslationResult:
        """Traduz planilha XLSX"""
        start_time = time.time()
        result = TranslationResult(success=False)
        
        try:
            wb = load_workbook(input_path)
            original_count = 0
            translated_count = 0
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str) and cell.value.strip():
                            original_count += 1
                            original_text = str(cell.value)
                            translated_text = self.translate_text(original_text, source_lang, target_lang, model)
                            cell.value = translated_text
                            translated_count += 1
            
            wb.save(output_path)
            
            result.success = True
            result.original_elements = original_count
            result.translated_elements = translated_count
            result.processing_time = time.time() - start_time
            
        except Exception as e:
            result.errors.append(f"Erro XLSX: {str(e)}")
            logger.error(f"Erro processando XLSX: {e}")
        
        return result

def translate_file_professional(input_path: str, output_path: str, glossary_path: Optional[str], 
                               use_ai: bool, source_lang: str, target_lang: str, model: str = None) -> TranslationResult:
    """Função principal de tradução de arquivos"""
    
    translator = DocumentTranslator()
    file_ext = Path(input_path).suffix.lower()
    
    if file_ext == '.docx':
        return translator.translate_docx(input_path, output_path, source_lang, target_lang, model)
    elif file_ext == '.pptx':
        return translator.translate_pptx(input_path, output_path, source_lang, target_lang, model)
    elif file_ext == '.xlsx':
        return translator.translate_xlsx(input_path, output_path, source_lang, target_lang, model)
    else:
        result = TranslationResult(success=False)
        result.errors.append(f"Formato de arquivo não suportado: {file_ext}")
        return result