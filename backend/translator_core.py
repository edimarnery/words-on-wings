# -*- coding: utf-8 -*-
"""
Sistema de Tradução Simplificado - Compatível com Docker
Inclui suporte a Espanhol Boliviano sem dependências externas
"""

import os
import json
import logging
import regex as re
from typing import Dict, List, Optional, Tuple
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation
from openai import OpenAI

DEFAULT_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")

# Configurar logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

_client = None
def get_client():
    global _client
    if _client is None and OPENAI_API_KEY:
        _client = OpenAI(api_key=OPENAI_API_KEY)
    return _client

# Mapeamento de idiomas melhorado com Espanhol Boliviano
IDIOMAS_MAPEAMENTO = {
    "Espanhol": "es",
    "Espanhol (Bolívia)": "es-bo",
    "Espanhol (Colômbia)": "es-co", 
    "Espanhol (México)": "es-mx",
    "Português do Brasil": "pt-br",
    "Português de Portugal": "pt-pt",
    "Inglês": "en",
    "Francês": "fr",
    "Italiano": "it",
    "Alemão": "de"
}

# Glossário básico integrado para documentos oficiais
GLOSSARIO_BASICO = {
    "CERTIDÃO DE NASCIMENTO": {
        "es": "CERTIFICADO DE NACIMIENTO",
        "es-bo": "CERTIFICADO DE NACIMIENTO"
    },
    "CERTIDÃO DE CASAMENTO": {
        "es": "CERTIFICADO DE MATRIMONIO", 
        "es-bo": "CERTIFICADO DE MATRIMONIO"
    },
    "CERTIDÃO DE ÓBITO": {
        "es": "CERTIFICADO DE DEFUNCIÓN",
        "es-bo": "CERTIFICADO DE DEFUNCIÓN"
    },
    "REGISTRO CIVIL DAS PESSOAS NATURAIS": {
        "es": "REGISTRO CIVIL DE LAS PERSONAS NATURALES",
        "es-bo": "REGISTRO CIVIL DE LAS PERSONAS NATURALES"
    },
    "REPÚBLICA FEDERATIVA DO BRASIL": {
        "es": "REPÚBLICA FEDERATIVA DE BRASIL",
        "es-bo": "REPÚBLICA FEDERATIVA DE BRASIL"
    }
}

def load_glossario_xlsx(path: str) -> Dict[str, str]:
    """Carrega glossário de arquivo Excel"""
    subs: Dict[str, str] = {}
    if not path or not os.path.isfile(path):
        return subs
    
    try:
        wb = load_workbook(path)
        ws = wb.active
        for row in ws.iter_rows(min_row=1, values_only=True):
            if not row or len(row) < 2: 
                continue
            o = (row[0] or "").strip()
            d = (row[1] or "").strip()
            if o and d: 
                subs[o.lower()] = d
    except Exception as e:
        logger.error(f"Erro ao carregar glossário Excel: {e}")
    
    return subs

def _build_gloss_regex(subs: Dict[str,str]):
    """Constrói regex para substituições do glossário"""
    terms = sorted(subs.keys(), key=len, reverse=True)
    if not terms: 
        return None, {}
    
    pats = []
    m = {}
    
    for t in terms:
        patt = r"\b" + re.sub(r"\s+", r"\\s+", re.escape(t)) + r"\b"
        gid = f"g{len(m)}"
        pats.append(f"(?P<{gid}>{patt})")
        m[gid] = subs[t]
    
    rx = re.compile("|".join(pats), flags=re.IGNORECASE|re.UNICODE)
    return rx, m

def _case_keep(src: str, dst: str) -> str:
    """Preserva o case da string original"""
    if src.isupper(): 
        return dst.upper()
    if src[:1].isupper() and (len(src)==1 or src[1:].islower()):
        return dst[:1].upper() + dst[1:]
    return dst

class Glossario:
    """Classe para gerenciar glossários de tradução"""
    
    def __init__(self, subs: Dict[str,str]):
        self.subs = subs
        self.rx, self.m = _build_gloss_regex(subs)
    
    def apply(self, text: str) -> str:
        """Aplica substituições do glossário"""
        if not self.rx:
            return text
        
        def repl_fn(match):
            for gid, replacement in self.m.items():
                if match.lastgroup == gid:
                    matched_text = match.group(gid)
                    return _case_keep(matched_text, replacement)
            return match.group(0)
        
        return self.rx.sub(repl_fn, text)

def aplicar_glossario_basico(texto: str, codigo_idioma: str) -> str:
    """Aplica glossário básico integrado"""
    texto_resultado = texto
    
    for termo_pt, traducoes in GLOSSARIO_BASICO.items():
        if codigo_idioma in traducoes:
            traducao = traducoes[codigo_idioma]
            # Substituir preservando case
            texto_resultado = re.sub(
                re.escape(termo_pt), 
                traducao, 
                texto_resultado, 
                flags=re.IGNORECASE
            )
    
    return texto_resultado

def detectar_codigo_idioma(idioma_nome: str) -> str:
    """Detecta código do idioma baseado no nome"""
    idioma_normalizado = idioma_nome.strip()
    
    if idioma_normalizado in IDIOMAS_MAPEAMENTO:
        return IDIOMAS_MAPEAMENTO[idioma_normalizado]
    
    idioma_lower = idioma_normalizado.lower()
    
    if "bolívia" in idioma_lower or "boliviano" in idioma_lower:
        return "es-bo"
    elif "espanhol" in idioma_lower or "español" in idioma_lower:
        return "es"
    elif "português" in idioma_lower and "brasil" in idioma_lower:
        return "pt-br"
    elif "inglês" in idioma_lower or "english" in idioma_lower:
        return "en"
    
    return "es"

def criar_prompt_especializado(texto: str, idioma_destino: str, codigo_idioma: str) -> str:
    """Cria prompt especializado para tradução"""
    
    prompt_base = f"""Você é um tradutor especializado em documentos oficiais brasileiros.

TAREFA: Traduzir o seguinte texto do português brasileiro para {idioma_destino}.

INSTRUÇÕES ESPECÍFICAS:
1. Mantenha EXATAMENTE a formatação original
2. Preserve TODOS os números, códigos e referências
3. Mantenha nomes próprios inalterados
4. Use terminologia jurídica apropriada
5. Preserve quebras de linha e espaçamento"""

    if codigo_idioma == "es-bo":
        prompt_base += """
6. Use espanhol formal apropriado para documentos oficiais bolivianos
7. Mantenha tratamento respeitoso e formal
8. Use terminologia jurídica padrão do espanhol"""
    elif codigo_idioma.startswith("es"):
        prompt_base += """
6. Use espanhol formal para documentos oficiais
7. Mantenha terminologia jurídica padrão"""
    
    if "CERTIDÃO" in texto.upper():
        prompt_base += """

TERMOS ESPECÍFICOS IMPORTANTES:
- "CERTIDÃO DE NASCIMENTO" = "CERTIFICADO DE NACIMIENTO"
- "REGISTRO CIVIL DAS PESSOAS NATURAIS" = "REGISTRO CIVIL DE LAS PERSONAS NATURALES"
- "FILIAÇÃO" = "FILIACIÓN"
- "AVÓS" = "ABUELOS"
- "Dou fé" = "Doy fe"
- "Escrevente Autorizada" = "Escribiente Autorizada"
"""
    
    prompt_base += f"""

TEXTO A TRADUZIR:
{texto}

TRADUÇÃO:"""
    
    return prompt_base

def ia_batch_translate(chunks: List[str], origem: str, destino: str, model: Optional[str]=None) -> Dict[str,str]:
    """Traduz chunks de texto via IA"""
    client = get_client()
    if not client:
        raise Exception("Cliente OpenAI não configurado")
    
    codigo_idioma = detectar_codigo_idioma(destino)
    result = {}
    
    for i, chunk in enumerate(chunks):
        if not chunk.strip():
            result[f"chunk_{i}"] = chunk
            continue
            
        try:
            prompt = criar_prompt_especializado(chunk, destino, codigo_idioma)
            
            response = client.chat.completions.create(
                model=model or DEFAULT_MODEL,
                messages=[
                    {"role": "system", "content": "Você é um tradutor especializado em documentos oficiais."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.1,
                max_tokens=4000
            )
            
            traducao = response.choices[0].message.content.strip()
            
            # Aplicar glossário básico
            traducao = aplicar_glossario_basico(traducao, codigo_idioma)
            
            result[f"chunk_{i}"] = traducao
            
        except Exception as e:
            logger.error(f"Erro na tradução do chunk {i}: {e}")
            result[f"chunk_{i}"] = chunk
    
    return result

def translate_text(text: str, gloss: Optional[Glossario], usar_ia: bool, origem: str, destino: str, model: Optional[str]=None) -> str:
    """Traduz texto usando IA"""
    if not usar_ia or not text.strip():
        return text
    
    chunks = [text]
    result = ia_batch_translate(chunks, origem, destino, model)
    traducao = result.get("chunk_0", text)
    
    if gloss:
        traducao = gloss.apply(traducao)
    
    return traducao

def _docx_runs(p, gloss, usar_ia, origem, destino, model):
    """Processa runs de parágrafo DOCX"""
    for run in p.runs:
        if run.text.strip():
            run.text = translate_text(run.text, gloss, usar_ia, origem, destino, model)

def translate_docx(inp, outp, gloss, usar_ia, origem, destino, model=None):
    """Traduz documento DOCX"""
    doc = Document(inp)
    for p in doc.paragraphs:
        _docx_runs(p, gloss, usar_ia, origem, destino, model)
    
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    _docx_runs(p, gloss, usar_ia, origem, destino, model)
    
    doc.save(outp)

def _pptx_tf(tf, gloss, usar_ia, origem, destino, model):
    """Processa text frame PPTX"""
    for p in tf.paragraphs:
        for run in p.runs:
            if run.text.strip():
                run.text = translate_text(run.text, gloss, usar_ia, origem, destino, model)

def translate_pptx(inp, outp, gloss, usar_ia, origem, destino, model=None):
    """Traduz apresentação PPTX"""
    prs = Presentation(inp)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                _pptx_tf(shape.text_frame, gloss, usar_ia, origem, destino, model)
    prs.save(outp)

def translate_xlsx(inp, outp, gloss, usar_ia, origem, destino, model=None):
    """Traduz planilha XLSX"""
    wb = load_workbook(inp)
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if cell.value and isinstance(cell.value, str) and cell.value.strip():
                    cell.value = translate_text(cell.value, gloss, usar_ia, origem, destino, model)
    wb.save(outp)

def translate_file(inp, outp, gloss_path, usar_ia, origem, destino, model=None):
    """Traduz arquivo - função principal de compatibilidade"""
    gloss = Glossario(load_glossario_xlsx(gloss_path)) if gloss_path else None
    ext = os.path.splitext(inp)[1].lower()
    
    if ext == '.docx':
        translate_docx(inp, outp, gloss, usar_ia, origem, destino, model)
    elif ext == '.pptx':
        translate_pptx(inp, outp, gloss, usar_ia, origem, destino, model)
    elif ext == '.xlsx':
        translate_xlsx(inp, outp, gloss, usar_ia, origem, destino, model)
    else:
        raise ValueError(f'Extensão não suportada: {ext}')